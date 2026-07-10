#!/usr/bin/env python3
"""Free-form PPTX generator for the 'claude-cowork' Talk."""
import os, re
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

TALK = "/sessions/affectionate-quirky-wright/mnt/talksmith-mim/talks/claude-cowork"
OUT  = os.path.join(TALK, "output")
SRC  = os.path.join(OUT, "final.free-form.intermediate.md")
BASE = "/sessions/affectionate-quirky-wright/mnt/.remote-plugins/plugin_01Ch8z21oXJFxDuxQfd3mBxk/config/pptx-styles/free-form/base-template.pptx"
LOG  = os.path.join(OUT, ".layout-log.md")

CORAL = RGBColor(0xC9,0x5B,0x3C)
INK   = RGBColor(0x1F,0x1E,0x1E)
CARD  = RGBColor(0xF5,0xF1,0xEC)
WHITE = RGBColor(0xFF,0xFF,0xFF)
MUTED = RGBColor(0x5A,0x54,0x52)
LIGHT = RGBColor(0xF5,0xF1,0xEC)

SW, SH = 9144000, 5143500
IN = 914400
def inch(v): return int(v*IN)

COVER = {
    "{{PRESENTATION_TITLE}}": "AI Generativa para Biomedicina",
    "{{TALK_SUBTITLE}}":      "Claude Cowork para el día a día",
    "{{PRESENTER}}":          "Paulo Veiga, Docente de Universidad Austral",
    "{{DATE}}":               "Junio 2026",
}

log_lines = ["# Layout log — free-form render (regenerated)\n"]
def log(s): log_lines.append(s)

def parse():
    lines = open(SRC, encoding="utf-8").read().splitlines()
    agenda=[]; blocks=[]
    i=0; n=len(lines)
    while i<n:
        s=lines[i].strip()
        if s.startswith("<!--"):
            i+=1; continue
        if s.startswith("# ") and not s.startswith("## "):
            title=s[2:].strip()
            if title.lower()=="agenda":
                i+=1
                while i<n:
                    t=lines[i].strip()
                    if t.startswith("#") or t=="---": break
                    m=re.match(r"-\s+(.*)", t)
                    if m: agenda.append(m.group(1).strip())
                    i+=1
                continue
            blocks.append(("divider", title)); i+=1; continue
        if s.startswith("## "):
            slide={"title":s[3:].strip(),"rawbody":[],"notes":""}
            blocks.append(("slide", slide)); i+=1
            body=[]; notes=[]; innotes=False; incode=False
            while i<n:
                l=lines[i]; st=l.strip()
                if st.startswith("<!--"): i+=1; continue
                if not incode and (st.startswith("## ") or (st.startswith("# ") and not st.startswith("## "))):
                    break
                if st=="---" and not incode:
                    i+=1; break
                if st.startswith("### Notes"):
                    innotes=True; i+=1; continue
                if st.startswith("### Sources"):
                    innotes=False; i+=1
                    while i<n and lines[i].strip()!="---" and not lines[i].strip().startswith("#"):
                        i+=1
                    continue
                if st.startswith("```"):
                    incode=not incode
                    if incode: body.append(("codeopen",None))
                    i+=1; continue
                if innotes: notes.append(l)
                else: body.append(("code",l) if incode else ("raw",l))
                i+=1
            slide["rawbody"]=body
            slide["notes"]="\n".join(notes).strip()
            continue
        i+=1
    return agenda, blocks

def build_blocks(rawbody):
    blocks=[]; tbl=None; code=[]
    for kind,l in rawbody:
        if kind=="code": code.append(l); continue
        if kind=="codeopen": continue
        st=l.strip()
        if st=="": continue
        if st.startswith("|"):
            if re.fullmatch(r"[-:\s|]*", st): continue
            cells=[c.strip() for c in st.strip().strip("|").split("|")]
            if tbl is None:
                tbl=[]; blocks.append(("table",tbl))
            tbl.append(cells); continue
        else: tbl=None
        m=re.search(r"!\[[^\]]*\]\((images/[^)]+)\)", st)
        if m: blocks.append(("image", m.group(1))); continue
        mb=re.match(r"^(\s*)-\s+(.*)$", l)
        if mb:
            blocks.append(("bullet", len(mb.group(1))//2, mb.group(2).strip())); continue
        blocks.append(("para", st))
    if code: blocks.append(("code","\n".join(code)))
    return blocks

def blank_slide(prs):
    best=min(prs.slide_layouts, key=lambda L: len(L.placeholders))
    sl=prs.slides.add_slide(best)
    for ph in list(sl.placeholders):
        ph._element.getparent().remove(ph._element)
    return sl

def bg(slide, color):
    cSld=slide._element
    old=cSld.find(qn('p:bg'))
    if old is not None: cSld.remove(old)
    xml=('<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
         'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
         '<p:bgPr><a:solidFill><a:srgbClr val="%02X%02X%02X"/></a:solidFill>'
         '<a:effectLst/></p:bgPr></p:bg>' % (color[0],color[1],color[2]))
    cSld.insert(0, parse_xml(xml))

def rect(slide,x,y,w,h,fill=None,shape=MSO_SHAPE.RECTANGLE):
    sp=slide.shapes.add_shape(shape,x,y,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.fill.background(); sp.shadow.inherit=False
    return sp

def add_runs(p,text,size,color,bold=False,font="Arial"):
    parts=re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", text)
    any_run=False
    for part in parts:
        if part=="": continue
        b=bold; it=False; t=part
        if part.startswith("**") and part.endswith("**"):
            t=part[2:-2]; b=True
        elif len(part)>2 and part.startswith("*") and part.endswith("*"):
            t=part[1:-1]; it=True
        r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.name=font
        r.font.bold=b; r.font.italic=it; r.font.color.rgb=color; any_run=True
    if not any_run:
        r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.name=font
        r.font.color.rgb=color; r.font.bold=bold

def textbox(slide,x,y,w,h,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Emu(0); tf.margin_right=Emu(0)
    tf.margin_top=Emu(0); tf.margin_bottom=Emu(0)
    return tb,tf

def set_noaspect(pic):
    cNvPicPr=pic._element.find(qn('p:nvPicPr')).find(qn('p:cNvPicPr'))
    if cNvPicPr.find(qn('a:picLocks')) is None:
        cNvPicPr.insert(0, cNvPicPr.makeelement(qn('a:picLocks'),{'noChangeAspect':'1'}))

def place_image(slide,relpath,bx,by,bw,bh):
    path=os.path.join(TALK,relpath)
    im=Image.open(path); nw,nh=im.size
    if nw/nh>=bw/bh:
        pic=slide.shapes.add_picture(path,bx,by,width=bw); pic.top=by+(bh-pic.height)//2
    else:
        pic=slide.shapes.add_picture(path,bx,by,height=bh); pic.left=bx+(bw-pic.width)//2
    set_noaspect(pic); return pic

def content_title(slide,title,sec_no,sec_name):
    rect(slide, inch(0.5), inch(0.34), inch(0.13), inch(0.62), fill=CORAL)
    tb,tf=textbox(slide, inch(0.75), inch(0.30), inch(8.7), inch(0.80), MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    add_runs(p,title,20,INK,bold=True)
    for r in p.runs: r.font.size=Pt(20)
    if sec_name:
        tb2,tf2=textbox(slide, inch(0.75), inch(1.03), inch(8.6), inch(0.24))
        p2=tf2.paragraphs[0]; r=p2.add_run()
        r.text="%02d · %s"%(sec_no, sec_name.upper())
        r.font.size=Pt(8); r.font.name="Arial"; r.font.bold=True; r.font.color.rgb=CORAL

def render_bullets(slide,blocks,x,y,w,h,size):
    tb,tf=textbox(slide,x,y,w,h); first=True
    for blk in blocks:
        if blk[0]=="bullet":
            _,level,text=blk
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            p.alignment=PP_ALIGN.LEFT; p.space_after=Pt(size*0.35); p.space_before=Pt(0)
            p.line_spacing=1.04
            pPr=p._p.get_or_add_pPr()
            pPr.set('marL', str(inch(level*0.28))); pPr.set('indent', str(-inch(0.22)))
            rm=p.add_run(); rm.text=("▪  " if level==0 else "–  ")
            rm.font.size=Pt(size); rm.font.name="Arial"; rm.font.bold=True; rm.font.color.rgb=CORAL
            add_runs(p,text, size if level==0 else size-0.5, INK if level==0 else MUTED)
        elif blk[0]=="para":
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            p.space_after=Pt(size*0.35); p.line_spacing=1.04
            txt=blk[1]; italic=txt.startswith("*") and txt.endswith("*")
            if italic: txt=txt.strip("*")
            add_runs(p,txt,size-0.5,MUTED)
            if italic:
                for r in p.runs: r.font.italic=True
    return tb

def render_table(slide,rows,x,y,w,maxh):
    nrows=len(rows); ncols=max(len(r) for r in rows)
    for r in rows:
        while len(r)<ncols: r.append("")
    rowh=int(min(maxh/nrows, inch(0.42)))
    gr=slide.shapes.add_table(nrows,ncols,x,y,w,rowh*nrows).table
    gr.first_row=False; gr.horz_banding=False
    for ci in range(ncols): gr.columns[ci].width=int(w/ncols)
    for ri,row in enumerate(rows):
        for ci in range(ncols):
            cell=gr.cell(ri,ci)
            cell.margin_left=Emu(inch(0.06)); cell.margin_right=Emu(inch(0.06))
            cell.margin_top=Emu(inch(0.02)); cell.margin_bottom=Emu(inch(0.02))
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = CORAL if ri==0 else (CARD if ri%2 else WHITE)
            p=cell.text_frame.paragraphs[0]; cell.text_frame.word_wrap=True
            r=p.add_run(); r.text=row[ci]; r.font.size=Pt(10.5); r.font.name="Arial"
            if ri==0: r.font.bold=True; r.font.color.rgb=WHITE
            else:
                r.font.color.rgb=INK
                if ci==0: r.font.bold=True
    return rowh*nrows

def build_content(prs,sd,sec_no,sec_name):
    slide=blank_slide(prs); bg(slide,WHITE)
    title=sd["title"]; blocks=build_blocks(sd["rawbody"])
    content_title(slide,title,sec_no,sec_name)
    images=[b for b in blocks if b[0]=="image"]
    tables=[b for b in blocks if b[0]=="table"]
    bullets=[b for b in blocks if b[0] in ("bullet","para")]
    nb=len(bullets); nimg=len(images); has_tbl=bool(tables)
    if nb<=4: fs=12.5
    elif nb<=6: fs=11.5
    elif nb<=8: fs=10.5
    elif nb<=10: fs=9.5
    else: fs=9.0
    body_top=inch(1.35); body_h=inch(5.42)-body_top
    tmpl="content+image" if nimg>=1 else "content-text"
    if has_tbl: tmpl="comparison(table)"
    if nimg==0 and not has_tbl:
        render_bullets(slide,bullets,inch(0.75),body_top,inch(8.7),body_h,fs)
    elif has_tbl:
        tx=inch(0.75); tw=inch(5.3) if nimg else inch(8.7)
        render_bullets(slide,bullets,tx,body_top,tw,inch(2.0),9.5)
        render_table(slide,tables[0][1],tx,body_top+inch(2.2),tw,inch(1.85))
        if nimg: place_image(slide,images[0][1],inch(6.25),inch(2.7),inch(3.25),inch(2.4))
    else:
        tx=inch(0.75); tw=inch(4.95); ix=inch(5.95); iw=inch(3.55)
        render_bullets(slide,bullets,tx,body_top,tw,body_h,fs)
        if nimg==1:
            place_image(slide,images[0][1],ix,body_top,iw,body_h)
        else:
            half=(body_h-inch(0.2))//2
            place_image(slide,images[0][1],ix,body_top,iw,half)
            place_image(slide,images[1][1],ix,body_top+half+inch(0.2),iw,half)
    if sd["notes"]: slide.notes_slide.notes_text_frame.text=sd["notes"]
    log("- «%s» — %s; %dimg %dtbl %dblk @%.1fpt"%(title[:44],tmpl,nimg,1 if has_tbl else 0,nb,fs))

def build_demo(prs,sd,sec_no,sec_name):
    slide=blank_slide(prs); bg(slide,WHITE)
    title=sd["title"]; blocks=build_blocks(sd["rawbody"])
    content_title(slide,title,sec_no,sec_name)
    images=[b for b in blocks if b[0]=="image"]
    bullets=[b for b in blocks if b[0] in ("bullet","para")]
    tx=inch(0.75); tw=inch(4.95)
    banner=rect(slide,tx,inch(1.35),tw,inch(0.6),fill=CORAL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf=banner.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text="▶  D E M O   T I M E"
    r.font.size=Pt(15); r.font.name="Arial"; r.font.bold=True; r.font.color.rgb=WHITE
    render_bullets(slide,bullets,tx,inch(2.15),tw,inch(3.2),10.5)
    if images: place_image(slide,images[0][1],inch(5.95),inch(1.35),inch(3.55),inch(4.0))
    if sd["notes"]: slide.notes_slide.notes_text_frame.text=sd["notes"]
    log("- «%s» — content+image (demo banner); 1img 0tbl %dblk @10.5pt"%(title[:44],len(bullets)))

def build_divider(prs,title,idx):
    slide=blank_slide(prs); bg(slide,INK)
    tb,tf=textbox(slide,inch(0.6),inch(0.7),inch(5.0),inch(2.2))
    p=tf.paragraphs[0]; r=p.add_run(); r.text="%02d"%idx
    r.font.size=Pt(120); r.font.name="Arial"; r.font.bold=True; r.font.color.rgb=CORAL
    rect(slide,inch(0.66),inch(3.05),inch(1.6),inch(0.10),fill=CORAL)
    tb2,tf2=textbox(slide,inch(0.66),inch(3.35),inch(8.6),inch(1.4))
    p2=tf2.paragraphs[0]; add_runs(p2,title,30,LIGHT,bold=True)
    for r in p2.runs: r.font.size=Pt(30)
    log("- Divider «%s» — dark, index %02d"%(title[:50],idx))

def build_agenda(prs,agenda):
    slide=blank_slide(prs); bg(slide,WHITE)
    rect(slide,inch(0.5),inch(0.42),inch(0.13),inch(0.62),fill=CORAL)
    tb,tf=textbox(slide,inch(0.75),inch(0.40),inch(8.6),inch(0.75),MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; r=p.add_run(); r.text="Agenda"
    r.font.size=Pt(26); r.font.name="Arial"; r.font.bold=True; r.font.color.rgb=INK
    top=inch(1.45); stride=inch(0.58)
    for k,item in enumerate(agenda):
        y=top+k*stride
        el=slide.shapes.add_shape(MSO_SHAPE.OVAL,inch(0.8),y,inch(0.42),inch(0.42))
        el.fill.solid(); el.fill.fore_color.rgb=CORAL; el.line.fill.background(); el.shadow.inherit=False
        ep=el.text_frame.paragraphs[0]; el.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
        ep.alignment=PP_ALIGN.CENTER
        m=re.match(r"(\d+)\.\s*(.*)", item)
        num=m.group(1) if m else str(k+1); label=m.group(2) if m else item
        er=ep.add_run(); er.text=num; er.font.size=Pt(14); er.font.name="Arial"
        er.font.bold=True; er.font.color.rgb=WHITE
        tb2,tf2=textbox(slide,inch(1.45),y,inch(7.9),inch(0.42),MSO_ANCHOR.MIDDLE)
        p2=tf2.paragraphs[0]; rr=p2.add_run(); rr.text=label
        rr.font.size=Pt(14.5); rr.font.name="Arial"; rr.font.color.rgb=INK
    log("- Agenda — %d sections, coral number chips"%len(agenda))

def fill_cover(prs):
    for sh in prs.slides[0].shapes:
        if not sh.has_text_frame: continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                for tok,val in COVER.items():
                    if tok in run.text: run.text=run.text.replace(tok,val)
    log("- Cover — base-template slide 1, 4 tokens substituted")

def main():
    agenda,blocks=parse()
    prs=Presentation(BASE); prs.slide_width=SW; prs.slide_height=SH
    fill_cover(prs); build_agenda(prs,agenda)
    div_idx=0; sec_no=0; sec_name=""
    for kind,payload in blocks:
        if kind=="divider":
            div_idx+=1; sec_no=div_idx; sec_name=payload
            build_divider(prs,payload,div_idx)
        else:
            if payload["title"].startswith("(Demo time)"):
                build_demo(prs,payload,sec_no,sec_name)
            else:
                build_content(prs,payload,sec_no,sec_name)
    out=os.path.join(OUT,"final.free-form.pptx"); prs.save(out)
    import shutil; shutil.copy(out, os.path.join(OUT,"final.pptx"))
    open(LOG,"w",encoding="utf-8").write("\n".join(log_lines)+"\n")
    print("slides:", len(prs.slides._sldIdLst))

if __name__=="__main__": main()
